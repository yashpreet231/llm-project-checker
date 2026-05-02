"""Extract text from uploaded files (PDF, DOCX, PPTX, TXT, MD, CSV)."""
from __future__ import annotations

import csv
import io
import logging

logger = logging.getLogger(__name__)


def extract_text(filename: str, content: bytes) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext in ("txt", "md", "markdown", "text", "log", "py", "js", "java", "c", "cpp", "h", "css", "html", "json", "xml", "yaml", "yml", "sh", "sql", "r"):
        return content.decode("utf-8", errors="replace")

    if ext == "csv":
        text = content.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        return "\n".join(" | ".join(row) for row in rows)

    if ext == "pdf":
        return _extract_pdf(content)

    if ext in ("docx", "doc"):
        return _extract_docx(content)

    if ext in ("pptx", "ppt"):
        return _extract_pptx(content)

    return content.decode("utf-8", errors="replace")


def _extract_pdf(content: bytes) -> str:
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                text = page.extract_text(x_tolerance=2, y_tolerance=3)
                if text:
                    pages.append(text)
        if pages:
            return "\n\n".join(pages)
    except Exception as e:
        logger.warning(f"pdfplumber extraction failed, falling back to PyPDF2: {e}")

    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(io.BytesIO(content))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages) if pages else "[Could not extract text from this PDF]"
    except Exception as e:
        logger.warning(f"PDF extraction failed: {e}")
        return "[Could not extract text from this PDF]"


def _extract_docx(content: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        logger.warning(f"DOCX extraction failed: {e}")
        return "[Could not extract text from this document]"


def _extract_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        logger.warning(f"PPTX extraction failed: {e}")
        return "[Could not extract text from this presentation]"
